# app.py

import os
from flask import Flask, request, render_template, redirect, url_for
from dotenv import load_dotenv

from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document

from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_core.prompts import PromptTemplate
from langchain_core.runnables import RunnablePassthrough
from langchain_core.output_parsers import StrOutputParser
from langchain_community.llms import Cohere

# Load API keys
load_dotenv()

UPLOAD_FOLDER = "uploads"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

app = Flask(__name__)
app.config["UPLOAD_FOLDER"] = UPLOAD_FOLDER

# Global state
embeddings_model = HuggingFaceEmbeddings(model_name='sentence-transformers/all-MiniLM-L6-v2')
vectorstore = None
retriever = None
rag_chain = None

# -------- Document Handlers --------
def extract_text(file_path):
    if file_path.endswith(".pdf"):
        reader = PdfReader(file_path)
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    elif file_path.endswith(".pptx"):
        ppt = Presentation(file_path)
        return "\n".join(shape.text for slide in ppt.slides for shape in slide.shapes if hasattr(shape, "text"))
    elif file_path.endswith(".docx"):
        doc = Document(file_path)
        return "\n".join(para.text for para in doc.paragraphs)
    return ""

def chunk_text(text):
    splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
    return splitter.split_text(text)

def build_rag_chain(retriever):
    prompt_template = """Answer the question as precisely as possible using the provided context.
If the answer is not contained in the context, say "answer not available in context".

Context:
{context}

Question:
{question}

Answer:"""
    prompt = PromptTemplate.from_template(prompt_template)
    llm = Cohere(
        model="command",
        temperature=0.1,
        cohere_api_key=os.getenv("cohere_api_key")
    )
    return (
        {"context": retriever | (lambda docs: "\n\n".join(d.page_content for d in docs)),
         "question": RunnablePassthrough()}
        | prompt
        | llm
        | StrOutputParser()
    )

# -------- Routes --------
@app.route("/", methods=["GET", "POST"])
def index():
    global vectorstore, retriever, rag_chain

    response = None

    if request.method == "POST":
        if "file" in request.files:
            uploaded_file = request.files["file"]
            if uploaded_file.filename:
                file_path = os.path.join(app.config["UPLOAD_FOLDER"], uploaded_file.filename)
                uploaded_file.save(file_path)
                text = extract_text(file_path)
                chunks = chunk_text(text)
                if vectorstore is None:
                    vectorstore = FAISS.from_texts(chunks, embedding=embeddings_model)
                    retriever = vectorstore.as_retriever(search_type="similarity", search_kwargs={"k": 6})
                    rag_chain = build_rag_chain(retriever)
                else:
                    vectorstore.add_texts(chunks)
                response = f"Uploaded and processed {uploaded_file.filename}"

        elif "question" in request.form:
            question = request.form["question"]
            if rag_chain:
                response = rag_chain.invoke(question)
            else:
                response = "Please upload a document first."

    return render_template("index.html", response=response)

# -------- Entry Point --------
if __name__ == "__main__":
    app.run(debug=False)  # 👈 disable auto-reloader that resets memory

